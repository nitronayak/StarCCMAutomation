import tkinter as tk
from tkinter import filedialog
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Open folder selection dialog
root = tk.Tk()
root.withdraw()
folder_path = filedialog.askdirectory()

# Get the folder name from the selected path
folder_name = os.path.basename(folder_path)

# Generate presentation using the images in the "Plots" and "Scenes" folders
presentation = Presentation()

# Add the title slide
slide_layout = presentation.slide_layouts[0]  # Title slide layout
slide = presentation.slides.add_slide(slide_layout)

title = slide.shapes.title
title.text = "49ers Racing Post-Processing Data"
title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title.text_frame.paragraphs[0].font.size = Pt(40)

subtitle = slide.placeholders[1]
subtitle.text = folder_name
subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add the "Powered by Nayak" text box
powered_by_text = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(0.5))
powered_by_tf = powered_by_text.text_frame
powered_by_p = powered_by_tf.add_paragraph()
powered_by_p.text = "Powered by Nayak"
powered_by_p.font.size = Pt(12)
powered_by_p.font.bold = False
powered_by_p.alignment = PP_ALIGN.CENTER

# Get the "Plots" folder path
plots_folder = os.path.join(folder_path, "Plots")

# Get the "Scenes" folder path
scenes_folder = os.path.join(folder_path, "Scenes")

# Retrieve all the image files from the "Plots" folder
plot_files = [file for file in os.listdir(plots_folder) if file.endswith('.jpg')]

# Retrieve all the image files from the "Scenes" folder
scene_files = [file for file in os.listdir(scenes_folder) if file.endswith('.jpg')]

# Sort the image files alphabetically
plot_files.sort()
scene_files.sort()

# Add each image from the "Plots" folder to a slide in the presentation
for plot_file in plot_files:
    slide_layout = presentation.slide_layouts[1]  # Use layout for content with caption
    slide = presentation.slides.add_slide(slide_layout)
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    image_path = os.path.join(plots_folder, plot_file)
    image = slide.shapes.add_picture(image_path, Inches(0), Inches(0))

    # Calculate the scaled width and height while maintaining aspect ratio
    max_width = slide_width - Inches(1)  # Subtract a margin
    max_height = slide_height - Inches(1.5)  # Subtract a margin and extra space for image name

    original_width = image.width
    original_height = image.height

    width_ratio = max_width / original_width
    height_ratio = max_height / original_height
    scale_ratio = min(width_ratio, height_ratio)

    scaled_width = int(original_width * scale_ratio)
    scaled_height = int(original_height * scale_ratio)

    # Calculate the left and top coordinates for centering the image
    left = int((slide_width - scaled_width) / 2)
    top = int((slide_height - scaled_height) / 2)

    # Set the position and size of the image
    image.left = left
    image.top = top
    image.width = scaled_width
    image.height = scaled_height

    # Add file name as the title of the slide
    file_name = os.path.splitext(plot_file)[0]
    slide.shapes.title.text = file_name
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(16)

    # Add file name as a label at the bottom of the image
    txBox = slide.shapes.add_textbox(Inches(0), slide_height - Inches(0.75), slide_width, Inches(0.75))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = file_name
    p.font.size = Pt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Add each image from the "Scenes" folder to a slide in the presentation
for scene_file in scene_files:
    slide_layout = presentation.slide_layouts[1]  # Use layout for content with caption
    slide = presentation.slides.add_slide(slide_layout)
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height

    image_path = os.path.join(scenes_folder, scene_file)
    image = slide.shapes.add_picture(image_path, Inches(0), Inches(0))

    # Calculate the scaled width and height while maintaining aspect ratio
    max_width = slide_width - Inches(1)  # Subtract a margin
    max_height = slide_height - Inches(1.5)  # Subtract a margin and extra space for image name

    original_width = image.width
    original_height = image.height

    width_ratio = max_width / original_width
    height_ratio = max_height / original_height
    scale_ratio = min(width_ratio, height_ratio)

    scaled_width = int(original_width * scale_ratio)
    scaled_height = int(original_height * scale_ratio)

    # Calculate the left and top coordinates for centering the image
    left = int((slide_width - scaled_width) / 2)
    top = int((slide_height - scaled_height) / 2)

    # Set the position and size of the image
    image.left = left
    image.top = top
    image.width = scaled_width
    image.height = scaled_height

    # Add file name as the title of the slide
    file_name = os.path.splitext(scene_file)[0]
    slide.shapes.title.text = file_name
    slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(16)

    # Add file name as a label at the bottom of the image
    txBox = slide.shapes.add_textbox(Inches(0), slide_height - Inches(0.75), slide_width, Inches(0.75))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.text = file_name
    p.font.size = Pt(12)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Save the presentation using the folder name as the file name
output_path = os.path.join(folder_path, f"{folder_name}.pptx")
presentation.save(output_path)

print("Presentation created successfully!")
